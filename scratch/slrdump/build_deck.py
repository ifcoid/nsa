# -*- coding: utf-8 -*-
"""Build the proposal PowerPoint deck from SLR session data."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
CHARTS = os.path.join(HERE, "charts")
OUT_DIR = os.path.join(HERE, "..", "..", "outputs")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "Proposal_Disertasi_Neural_Decoding_BCI.pptx")

# ---- Palette ----
NAVY   = RGBColor(0x0B, 0x25, 0x45)
NAVY2  = RGBColor(0x13, 0x36, 0x5E)
TEAL   = RGBColor(0x13, 0xA0, 0xA6)
ACCENT = RGBColor(0xF4, 0xA3, 0x00)
GREY   = RGBColor(0x5B, 0x6B, 0x7B)
LGREY  = RGBColor(0xEE, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xC4, 0x45, 0x36)
GREEN  = RGBColor(0x2E, 0x7D, 0x5B)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for seg in para:
            text, size, color, bold, italic = (list(seg) + [False, False])[:5]
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = FONT
    return tb


def bg(s, color=WHITE):
    rect(s, 0, 0, SW, SH, fill=color)


def header(s, kicker, title, num):
    """Standard content-slide header. Returns y where body can start."""
    bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(1.18), fill=NAVY)
    rect(s, 0, Inches(1.18), SW, Pt(3), fill=TEAL)
    rect(s, 0, 0, Inches(0.16), Inches(1.18), fill=ACCENT)
    txt(s, Inches(0.55), Inches(0.16), Inches(11), Inches(0.3),
        [(kicker.upper(), 11, TEAL, True, False)])
    txt(s, Inches(0.55), Inches(0.44), Inches(11.8), Inches(0.7),
        [(title, 23, WHITE, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    # slide number badge
    txt(s, Inches(12.5), Inches(0.16), Inches(0.7), Inches(0.3),
        [(num, 11, GREY, False, False)], align=PP_ALIGN.RIGHT)
    return Inches(1.45)


def footer(s):
    txt(s, Inches(0.55), Inches(7.12), Inches(9), Inches(0.3),
        [[("Proposal Disertasi  ·  ", 8, GREY, False, False),
          ("Neural Decoding BCI berbasis State-Space Models", 8, GREY, False, True)]])
    txt(s, Inches(11.0), Inches(7.12), Inches(1.8), Inches(0.3),
        [("Rolly M. Awangga · Tel-U", 8, GREY, False, False)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, text, color, tcolor=WHITE, size=10, h=Inches(0.34)):
    r = rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    r.adjustments[0] = 0.5
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = text
    rn.font.size = Pt(size); rn.font.bold = True; rn.font.color.rgb = tcolor; rn.font.name = FONT
    return r


def card(s, x, y, w, h, fill=LGREY, line=None):
    r = rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             line=line, line_w=Pt(1.25) if line else None)
    r.adjustments[0] = 0.06
    return r


def bullets(s, x, y, w, h, items, size=14, color=NAVY, gap=7, bullet_color=TEAL,
            line_spacing=1.02):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else 3)
        p.line_spacing = line_spacing
        p.level = lvl
        r = p.add_run()
        r.text = ("●  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(size if lvl == 0 else size-2)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = FONT
        # color the bullet glyph
        r.font.bold = False
    return tb


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide()
bg(s, NAVY)
# decorative bands
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(5.4), SW, Inches(2.1), fill=NAVY2)
rect(s, 0, Inches(5.4), SW, Pt(3), fill=TEAL)
# accent neural dots motif (top-right)
import random
random.seed(7)
for _ in range(26):
    cx = Inches(9.2 + random.random()*3.8)
    cy = Inches(0.3 + random.random()*2.4)
    d = Pt(4 + random.random()*7)
    c = TEAL if random.random() > 0.5 else ACCENT
    o = rect(s, cx, cy, d, d, fill=c, shape=MSO_SHAPE.OVAL)
chip(s, Inches(0.8), Inches(0.85), Inches(3.0), "PROPOSAL DISERTASI", ACCENT, NAVY, 12)
txt(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.6),
    [[("Peningkatan Akurasi dan Skalabilitas ", 40, WHITE, True, False),
      ("Neural Decoding", 40, TEAL, True, False)],
     [("Menggunakan Model Machine Learning Lanjut", 40, WHITE, True, False)],
     [("pada Ekosistem BCI", 40, WHITE, True, False)]],
    line_spacing=1.02)
txt(s, Inches(0.82), Inches(4.35), Inches(11.5), Inches(0.8),
    [[("Tinjauan Sistematis Transisi Arsitektur Transformer → ", 16, GREY, False, True),
      ("State-Space Models (Mamba/S4)", 16, ACCENT, False, True)]])
txt(s, Inches(0.8), Inches(5.6), Inches(8.5), Inches(1.6),
    [[("Rolly Maulana Awangga", 18, WHITE, True, False),
      ("   NIM 303012510004", 13, GREY, False, False)],
     [("Program Doktor — Universitas Telkom", 13, GREY, False, False)],
     [("rollymaulanaa@student.telkomuniversity.ac.id", 12, TEAL, False, False)]], space_after=3)
txt(s, Inches(10.0), Inches(6.5), Inches(2.7), Inches(0.5),
    [[("Dipresentasikan: Juni 2026", 11, GREY, False, True)]], align=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 2 — KONTEKS & LATAR BELAKANG
# =====================================================================
s = slide()
y = header(s, "Latar Belakang", "Konteks: Neural Decoding pada Brain–Computer Interface", "2")
txt(s, Inches(0.55), y, Inches(12.2), Inches(0.9),
    [[("Brain–Computer Interface (BCI) menerjemahkan ", 15, NAVY, False, False),
      ("sinyal saraf", 15, TEAL, True, False),
      (" (EEG, fMRI, ECoG, intracortical spiking) menjadi perintah kontrol. "
       "Inti sistemnya adalah ", 15, NAVY, False, False),
      ("neural decoder", 15, TEAL, True, False),
      (" — model machine learning yang memetakan aktivitas otak ke niat pengguna.",
       15, NAVY, False, False)]], line_spacing=1.05)

cards = [
    ("⚡ Akurasi", "Deep learning (CNN, RNN, Transformer) telah mendongkrak akurasi dekoding secara signifikan dalam dekade terakhir.", TEAL),
    ("⏱️ Latensi Real-time", "BCI closed-loop & neuro-prostetik menuntut inferensi seketika (<10 ms) di perangkat edge berdaya rendah.", ACCENT),
    ("📈 Skalabilitas", "Sekuens sinyal saraf sangat panjang; biaya komputasi model harus tumbuh efisien, bukan kuadratik.", NAVY),
]
cw = Inches(3.95); gap = Inches(0.22); x0 = Inches(0.55); cy = Inches(2.7); ch = Inches(2.5)
for i, (t, d, col) in enumerate(cards):
    x = x0 + i*(cw+gap)
    card(s, x, cy, cw, ch, fill=LGREY)
    rect(s, x, cy, cw, Pt(5), fill=col)
    txt(s, x+Inches(0.25), cy+Inches(0.28), cw-Inches(0.5), Inches(0.6),
        [(t, 18, col, True, False)])
    txt(s, x+Inches(0.25), cy+Inches(1.0), cw-Inches(0.5), Inches(1.3),
        [(d, 13, GREY, False, False)], line_spacing=1.06)

card(s, Inches(0.55), Inches(5.5), Inches(12.23), Inches(1.35), fill=NAVY)
txt(s, Inches(0.85), Inches(5.62), Inches(11.6), Inches(1.1),
    [[("Tensi utama bidang ini:  ", 15, ACCENT, True, False),
      ("memaksimalkan akurasi dekoding sekaligus menekan biaya komputasi & latensi "
       "agar BCI dapat dipakai nyata di perangkat implan/edge.", 15, WHITE, False, False)],
     [("Pertanyaannya: arsitektur ML mana yang paling 'mangkus' untuk ekosistem BCI generasi berikutnya?",
       14, GREY, False, True)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s)

# =====================================================================
# SLIDE 3 — MASALAH & MOTIVASI
# =====================================================================
s = slide()
y = header(s, "Masalah", "Bottleneck Transformer & Munculnya State-Space Models", "3")
# Left: Transformer problem
card(s, Inches(0.55), y, Inches(5.9), Inches(4.0), fill=LGREY, line=RED)
chip(s, Inches(0.8), y+Inches(0.22), Inches(3.2), "ARSITEKTUR SAAT INI: TRANSFORMER", RED, WHITE, 11)
bullets(s, Inches(0.85), y+Inches(0.8), Inches(5.3), Inches(3.0), [
    "Akurasi superior pada dekoding BCI — sudah menjadi standar de-facto.",
    "Mekanisme atensi berkompleksitas KUADRATIK O(n²) terhadap panjang sekuens.",
    "Boros memori & daya → sulit ditanam pada implan / edge device.",
    "Gagal memenuhi target latensi <10 ms untuk closed-loop real-time.",
], size=13.5, gap=9)
# Right: SSM solution
card(s, Inches(6.85), y, Inches(5.93), Inches(4.0), fill=LGREY, line=TEAL)
chip(s, Inches(7.1), y+Inches(0.22), Inches(3.6), "PERGESERAN: STATE-SPACE MODELS (MAMBA/S4)", TEAL, WHITE, 11)
bullets(s, Inches(7.15), y+Inches(0.8), Inches(5.4), Inches(3.0), [
    "Kompleksitas LINEAR O(n) terhadap panjang sekuens.",
    "Hemat memori, mendukung paralelisasi layaknya konvolusi.",
    "Klaim latensi inferensi real-time turun tajam (mis. POSSM).",
    "Mewarisi akurasi Transformer tanpa beban komputasi kuadratik.",
], size=13.5, gap=9, bullet_color=TEAL)
# bottom strip
card(s, Inches(0.55), Inches(5.85), Inches(12.23), Inches(1.0), fill=NAVY)
txt(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.8),
    [[("Motivasi:  ", 15, ACCENT, True, False),
      ("Sejak S4 (2021) dan Mamba (2023), literatur deep-learning BCI bergeser cepat ke SSM — "
       "namun bukti efektivitasnya masih ", 14, WHITE, False, False),
      ("terfragmentasi dan belum disintesis", 14, ACCENT, True, False),
      (".", 14, WHITE, False, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.04)
footer(s)

# =====================================================================
# SLIDE 4 — RESEARCH GAP
# =====================================================================
s = slide()
y = header(s, "Kesenjangan Penelitian", "Research Gap — Fragmentasi Literatur (Tipe A)", "4")
chip(s, Inches(0.55), y, Inches(2.7), "GAP TIPE A · FRAGMENTASI", ACCENT, NAVY, 12)
txt(s, Inches(0.55), y+Inches(0.55), Inches(12.2), Inches(1.5),
    [[("Implementasi SSM/Mamba untuk neural decoding tersebar pada ", 16, NAVY, False, False),
      ("modalitas spesifik", 16, TEAL, True, False),
      (" (spiking NHP, EEG, fMRI) ", 16, NAVY, False, False),
      ("tanpa sintesis sistematis", 16, RED, True, False),
      (" mengenai rasio ", 16, NAVY, False, False),
      ("trade-off antara efisiensi komputasi riil (latensi <10 ms) dan akurasi generalisasi",
       16, NAVY, True, False),
      (" lintas task BCI.", 16, NAVY, False, False)]], line_spacing=1.1)

# three fragment islands + arrow to synthesis
frags = [("EEGMamba", "EEG"), ("POSSM", "Spiking NHP"), ("Energy-Guided\nMamba", "Topologi/EEG")]
fx = Inches(0.7); fy = Inches(3.4); fw = Inches(2.5); fh = Inches(1.5)
for i, (n, m) in enumerate(frags):
    x = fx + i*(fw+Inches(0.35))
    card(s, x, fy, fw, fh, fill=LGREY, line=GREY)
    txt(s, x, fy+Inches(0.25), fw, Inches(0.7),
        [(n, 15, NAVY, True, False)], align=PP_ALIGN.CENTER)
    txt(s, x, fy+Inches(0.95), fw, Inches(0.4),
        [("modalitas: "+m, 11, GREY, False, True)], align=PP_ALIGN.CENTER)
arrow = rect(s, Inches(9.0), fy+Inches(0.45), Inches(0.9), Inches(0.6),
             fill=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
card(s, Inches(10.05), fy-Inches(0.05), Inches(2.75), fh+Inches(0.1), fill=NAVY)
txt(s, Inches(10.1), fy+Inches(0.15), Inches(2.65), fh-Inches(0.2),
    [[("SINTESIS\nSISTEMATIS", 15, ACCENT, True, False)],
     [("(belum ada — kontribusi riset ini)", 11, WHITE, False, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

card(s, Inches(0.55), Inches(5.45), Inches(12.23), Inches(1.4), fill=LGREY, line=TEAL)
txt(s, Inches(0.85), Inches(5.6), Inches(11.6), Inches(1.1),
    [[("Bukti urgensi:  ", 14, TEAL, True, False),
      ("Arsitektur hibrida POSSM (Transformer–Mamba) mengklaim penurunan tajam waktu inferensi "
       "real-time, dan implementasi seperti 'Energy-Guided Brain Topology Mamba' (2025) "
       "menunjukkan tren berkembang pesat ", 13.5, NAVY, False, False),
      ("namun kelemahannya belum pernah disintesis.", 13.5, RED, True, False)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s)

# =====================================================================
# SLIDE 5 — PRIOR REVIEWS / NOVELTY
# =====================================================================
s = slide()
y = header(s, "Posisi terhadap Literatur", "Tinjauan Terdahulu & Klaim Kebaruan", "5")
rows = [
    ("MDPI (2025)", "Tren & limitasi riset BCI berbasis Transformer", "Stagnan di dalam framework Transformer; tak menyinggung transisi ke SSM/Mamba", "BEDA FRAMEWORK"),
    ("IntechOpen (2026)", "Biosensing → AI-powered decoding (multimodal)", "Cakupan terlalu luas; gagal mengagregasi efisiensi inferensi level mikro (SSM)", "BEDA FOKUS"),
    ("The Innovation (2026)", "Generative AI untuk dekoding BCI", "Fiksasi pada akurasi rekonstruksi (software); abai pada arsitektur boros daya di edge", "BEDA METODE"),
]
# table header
tx = Inches(0.55); tw = Inches(12.23)
colw = [Inches(2.5), Inches(3.5), Inches(4.4), Inches(1.83)]
hh = Inches(0.5)
heads = ["Tinjauan", "Cakupan", "Keterbatasan (Gap)", "Selisih"]
xx = tx
rect(s, tx, y, tw, hh, fill=NAVY)
for h, w in zip(heads, colw):
    txt(s, xx+Inches(0.1), y, w-Inches(0.2), hh, [(h, 12, WHITE, True, False)],
        anchor=MSO_ANCHOR.MIDDLE)
    xx += w
ry = y + hh
for ridx, row in enumerate(rows):
    rh = Inches(1.05)
    fillc = LGREY if ridx % 2 == 0 else WHITE
    rect(s, tx, ry, tw, rh, fill=fillc)
    xx = tx
    for ci, (cell, w) in enumerate(zip(row, colw)):
        if ci == 0:
            txt(s, xx+Inches(0.1), ry, w-Inches(0.2), rh, [(cell, 12, NAVY, True, False)],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        elif ci == 3:
            chip(s, xx+Inches(0.12), ry+Inches(0.32), w-Inches(0.24), cell, ACCENT, NAVY, 9.5)
        else:
            col = GREY if ci == 2 else NAVY
            txt(s, xx+Inches(0.1), ry, w-Inches(0.2), rh, [(cell, 11.5, col, False, ci==2)],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        xx += w
    ry += rh
card(s, tx, ry+Inches(0.12), tw, Inches(0.95), fill=TEAL)
txt(s, tx+Inches(0.3), ry+Inches(0.2), tw-Inches(0.6), Inches(0.8),
    [[("KEBARUAN:  ", 14, NAVY, True, False),
      ("Belum ada systematic review yang khusus mensintesis arsitektur SSM/Mamba untuk neural "
       "decoding BCI. Riset ini menjadi ", 13.5, WHITE, False, False),
      ("tinjauan sistematis komprehensif pertama", 13.5, NAVY, True, False),
      (" pada irisan tersebut.", 13.5, WHITE, False, False)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.04)
footer(s)

# =====================================================================
# SLIDE 6 — RESEARCH QUESTIONS
# =====================================================================
s = slide()
y = header(s, "Pertanyaan Penelitian", "Research Questions (1 Primary + 3 Secondary)", "6")
rqs = [
    ("RQ1", "PRIMARY", NAVY,
     "Sejauh mana arsitektur SSM (Mamba & hibrida) mencapai efisiensi 'mangkus' "
     "(latensi <10 ms) tanpa mengorbankan akurasi generalisasi lintas modalitas, "
     "dibanding arsitektur sebelumnya?"),
    ("RQ2", "SECONDARY", TEAL,
     "Bagaimana variasi performa akurasi model Mamba/S4 ketika diimplementasikan pada "
     "modalitas sinyal berbeda (spiking NHP, EEG, fMRI)?"),
    ("RQ3", "SECONDARY", TEAL,
     "Bagaimana konfigurasi arsitektur SSM (backbone murni vs hibrida Mamba-Transformer "
     "seperti POSSM) dan keunggulannya dibanding baseline konvensional?"),
    ("RQ4", "SECONDARY", TEAL,
     "Apa bukti empiris parameter 'mangkus' Mamba (memori, throughput, latensi <10 ms) "
     "untuk kelayakan neuro-prostetik real-time di edge devices?"),
]
cw = Inches(6.0); ch = Inches(2.45); gx = Inches(0.23); gy = Inches(0.25)
x0 = Inches(0.55); yy0 = y
for i, (tag, kind, col, q) in enumerate(rqs):
    r, c = divmod(i, 2)
    x = x0 + c*(cw+gx); yy = yy0 + r*(ch+gy)
    card(s, x, yy, cw, ch, fill=LGREY)
    rect(s, x, yy, Inches(0.14), ch, fill=col)
    txt(s, x+Inches(0.32), yy+Inches(0.2), Inches(2), Inches(0.5),
        [[(tag+"  ", 22, col, True, False)]])
    chip(s, x+Inches(1.85), yy+Inches(0.28), Inches(1.7), kind, col, WHITE, 10)
    txt(s, x+Inches(0.32), yy+Inches(0.95), cw-Inches(0.6), Inches(1.4),
        [(q, 13, NAVY, False, False)], line_spacing=1.06)
footer(s)

# =====================================================================
# SLIDE 7 — PICO FRAMEWORK
# =====================================================================
s = slide()
y = header(s, "Kerangka Konseptual", "Definisi PICO & Terminologi Kanonik", "7")
pico = [
    ("P", "Populasi", "Sinyal saraf pengguna BCI lintas modalitas (EEG, fMRI, NHP spiking) dalam tugas neural decoding", NAVY),
    ("I", "Intervensi", "Algoritma State-Space Models (Mamba, S4) dan arsitektur hibridanya", TEAL),
    ("C", "Komparasi", "Arsitektur deep learning sebelumnya (Transformer, CNN, RNN) atau tanpa pembanding", GREY),
    ("O", "Outcome", "Akurasi dekoding (WAJIB) + metrik efisiensi komputasi / latensi <10 ms (OPSIONAL)", ACCENT),
]
cw = Inches(2.95); gap = Inches(0.15); x0 = Inches(0.55); cy = y; ch = Inches(3.1)
for i, (l, t, d, col) in enumerate(pico):
    x = x0 + i*(cw+gap)
    card(s, x, cy, cw, ch, fill=LGREY)
    circ = rect(s, x+cw/2-Inches(0.45), cy+Inches(0.28), Inches(0.9), Inches(0.9),
                fill=col, shape=MSO_SHAPE.OVAL)
    txt(s, x+cw/2-Inches(0.45), cy+Inches(0.28), Inches(0.9), Inches(0.9),
        [(l, 30, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, cy+Inches(1.35), cw, Inches(0.4), [(t, 15, col, True, False)], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), cy+Inches(1.85), cw-Inches(0.4), Inches(1.2),
        [(d, 11.5, NAVY, False, False)], align=PP_ALIGN.CENTER, line_spacing=1.05)
card(s, Inches(0.55), Inches(4.95), Inches(12.23), Inches(1.85), fill=NAVY)
txt(s, Inches(0.85), Inches(5.1), Inches(11.6), Inches(0.5),
    [[("Terminologi Kanonik:  ", 14, ACCENT, True, False),
      ("Structured State-Space Sequence Models (SSMs)", 14, WHITE, True, False)]])
txt(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.2),
    [[("Kelas arsitektur deep learning pemodelan sekuensial dengan parameterisasi ruang-keadaan kontinu "
       "→ kompleksitas LINEAR terhadap sekuens panjang, jauh lebih hemat memori dibanding atensi kuadratik Transformer.",
       12.5, WHITE, False, False)],
     [("Alternatif ditolak: ", 11.5, GREY, True, True),
      ("Kalman Filter (kontrol linear klasik) & RNN (pemrosesan strictly sekuensial, vanishing gradient).",
       11.5, GREY, False, True)]], line_spacing=1.05, space_after=6)
footer(s)

# =====================================================================
# SLIDE 8 — SCOPE + FINER
# =====================================================================
s = slide()
y = header(s, "Ruang Lingkup & Kelayakan", "Scope Filters & Validasi FINER", "8")
# left: scope
card(s, Inches(0.55), y, Inches(5.8), Inches(4.7), fill=LGREY)
chip(s, Inches(0.8), y+Inches(0.22), Inches(2.6), "RUANG LINGKUP", NAVY, WHITE, 11)
scope = [
    ("Rentang Tahun", "2021–2027 (era S4 → Mamba)"),
    ("Geografis", "Global (tanpa restriksi)"),
    ("Sektor", "Computer Science · Biomedis · Neurosains"),
    ("Bahasa", "English"),
    ("Tipe Dokumen", "Jurnal peer-reviewed + Prosiding konferensi"),
]
sy = y+Inches(0.85)
for k, v in scope:
    rect(s, Inches(0.8), sy+Inches(0.05), Inches(0.12), Inches(0.55), fill=TEAL)
    txt(s, Inches(1.05), sy, Inches(5.1), Inches(0.6),
        [[(k+":  ", 13, NAVY, True, False), (v, 13, GREY, False, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    sy += Inches(0.78)
# right: FINER
card(s, Inches(6.55), y, Inches(6.23), Inches(4.7), fill=LGREY, line=GREEN)
chip(s, Inches(6.8), y+Inches(0.22), Inches(3.6), "VALIDASI FINER — LULUS (PASS)", GREEN, WHITE, 11)
finer = [
    ("F", "Feasible", "30–50+ literatur primer relevan (EEGMamba, POSSM, MSResBiMamba, dll.)"),
    ("I", "Interesting", "Krusial bagi pengembang hardware BCI & neuro-engineer"),
    ("N", "Novel", "Sintesis pertama khusus SSM/Mamba — menambal blind-spot review terdahulu"),
    ("E", "Ethical", "Risiko minim (studi literatur); patuh lisensi data sekunder"),
    ("R", "Relevant", "Mendukung SDG 3 (Kesehatan) & SDG 9 (Inovasi & Infrastruktur)"),
]
fy = y+Inches(0.85)
for l, t, d in finer:
    c = rect(s, Inches(6.8), fy, Inches(0.55), Inches(0.55), fill=GREEN, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(6.8), fy, Inches(0.55), Inches(0.55), [(l, 16, WHITE, True, False)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.5), fy-Inches(0.02), Inches(5.1), Inches(0.7),
        [[(t+":  ", 12.5, NAVY, True, False), (d, 12, GREY, False, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    fy += Inches(0.78)
footer(s)

# =====================================================================
# SLIDE 9 — METODOLOGI (Agentic AI pipeline)
# =====================================================================
s = slide()
y = header(s, "Metodologi", "Pipeline SLR Otomatis berbasis Agentic AI", "9")
txt(s, Inches(0.55), y, Inches(12.2), Inches(0.6),
    [[("SLR dijalankan oleh sistem ", 14, NAVY, False, False),
      ("agentic AI", 14, TEAL, True, False),
      (" buatan sendiri (Go) — mengikuti protokol PRISMA, dengan dua agen reviewer (strict & liberal) "
       "dan validasi reliabilitas Cohen's Kappa di setiap tahap.", 14, NAVY, False, False)]],
    line_spacing=1.05)
mods = [
    ("M1", "Foundation", "Topik & gap", True),
    ("M2", "PICO", "RQ · FINER", True),
    ("M3", "Search", "Database · string", True),
    ("M4", "Mining", "Dedup · audit", True),
    ("M5", "Screening", "Kappa · eksklusi", True),
    ("M6", "Acquisition", "Full-text · vektor", "current"),
    ("M7", "Extraction", "Ekstraksi data", False),
    ("M8", "Synthesis", "Sintesis · bibliometrik", False),
    ("M9", "Manuscript", "Penyusunan naskah", False),
]
cw = Inches(1.3); gap = Inches(0.06); x0 = Inches(0.55); cy = Inches(2.6); ch = Inches(2.0)
for i, (m, t, d, st) in enumerate(mods):
    x = x0 + i*(cw+gap)
    if st == "current":
        fillc = ACCENT; tcol = NAVY; sub = NAVY
    elif st:
        fillc = NAVY; tcol = WHITE; sub = TEAL
    else:
        fillc = LGREY; tcol = GREY; sub = GREY
    card(s, x, cy, cw, ch, fill=fillc)
    txt(s, x, cy+Inches(0.25), cw, Inches(0.5), [(m, 22, tcol, True, False)], align=PP_ALIGN.CENTER)
    txt(s, x, cy+Inches(0.85), cw, Inches(0.5), [(t, 11.5, tcol, True, False)],
        align=PP_ALIGN.CENTER, line_spacing=0.95)
    txt(s, x+Inches(0.05), cy+Inches(1.3), cw-Inches(0.1), Inches(0.6),
        [(d, 9, sub, False, False)], align=PP_ALIGN.CENTER, line_spacing=0.95)
    if i < len(mods)-1:
        rect(s, x+cw-Inches(0.02), cy+ch/2-Inches(0.06), gap+Inches(0.04), Pt(2.5), fill=GREY)
# legend / current marker
txt(s, Inches(0.55), Inches(4.8), Inches(12), Inches(0.4),
    [[("■ ", 12, NAVY, True, False), ("Selesai (M1–M5)   ", 12, GREY, False, False),
      ("■ ", 12, ACCENT, True, False), ("Sedang berjalan (M6 — posisi saat ini)   ", 12, GREY, True, False),
      ("■ ", 12, GREY, True, False), ("Rencana (M7–M9)", 12, GREY, False, False)]])
card(s, Inches(0.55), Inches(5.35), Inches(12.23), Inches(1.5), fill=LGREY, line=TEAL)
txt(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.2),
    [[("Keunggulan pendekatan:  ", 13.5, TEAL, True, False),
      ("(1) reprodusibel & terdokumentasi penuh di MongoDB;  (2) dual-reviewer otomatis dengan "
       "audit Kappa;  (3) full-text di-vektorisasi ke Qdrant untuk ekstraksi semantik (RAG) pada M7–M8.",
       13, NAVY, False, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)
footer(s)

# =====================================================================
# SLIDE 10 — SEARCH STRATEGY
# =====================================================================
s = slide()
y = header(s, "Strategi Pencarian (Modul 3)", "Multi-Database + Grey Literature", "10")
dbs = [
    ("Scopus", "Jangkar multidisiplin tervalidasi", NAVY),
    ("IEEE Xplore", "Prosiding konferensi BCI/teknik", TEAL),
    ("PubMed", "Translasi klinis neuro-prostetik", GREY),
    ("arXiv / Scholar", "Grey literature — preprint SSM terbaru", ACCENT),
]
cw = Inches(2.95); gap = Inches(0.15); x0 = Inches(0.55); cy = y
for i, (n, d, col) in enumerate(dbs):
    x = x0 + i*(cw+gap)
    card(s, x, cy, cw, Inches(1.5), fill=LGREY)
    rect(s, x, cy, cw, Pt(5), fill=col)
    txt(s, x+Inches(0.2), cy+Inches(0.22), cw-Inches(0.4), Inches(0.5),
        [(n, 15, col, True, False)], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), cy+Inches(0.72), cw-Inches(0.4), Inches(0.7),
        [(d, 11, GREY, False, False)], align=PP_ALIGN.CENTER, line_spacing=1.0)
txt(s, Inches(0.55), cy+Inches(1.7), Inches(12), Inches(0.4),
    [[("Keputusan: ", 13, NAVY, True, False),
      ("Multi-database + grey literature wajib — SSM/Mamba terlalu baru (akhir 2023), index tradisional "
       "punya publication-lag bias.", 12.5, GREY, False, True)]])
# search string box
card(s, Inches(0.55), cy+Inches(2.25), Inches(12.23), Inches(2.45), fill=NAVY)
txt(s, Inches(0.8), cy+Inches(2.38), Inches(11.7), Inches(0.4),
    [[("Search String Inti (Scopus) — dieksekusi 2026-05-25", 13, ACCENT, True, False)]])
qstr = ('TITLE-ABS-KEY( ("Brain-Computer Interface" OR "BCI" OR "neural decoding" OR "neural signal*" '
        'OR "EEG" OR "fMRI" OR "fNIRS" OR "ECoG" OR "intracortical spiking" ...) AND '
        '("Structured State-Space Sequence Model*" OR "State-Space Model*" OR "SSM" OR '
        '"Mamba" OR "S4" OR "S4D" OR "POSSM") )')
tb = s.shapes.add_textbox(Inches(0.8), cy+Inches(2.85), Inches(11.7), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = qstr
r.font.size = Pt(12.5); r.font.color.rgb = WHITE; r.font.name = "Consolas"
p.line_spacing = 1.15
txt(s, Inches(0.8), cy+Inches(4.25), Inches(11.7), Inches(0.4),
    [[("Pre-validasi: ", 11, GREY, True, True),
      ("sintaks AMAN · trap-keyword (Mamba/S4/SSM) tersaring oleh AND konsep saraf · zona Goldilocks.",
       11, GREY, False, True)]])
footer(s)

# =====================================================================
# SLIDE 11 — MINING RESULTS (year chart)
# =====================================================================
s = slide()
y = header(s, "Hasil Identifikasi (Modul 4)", "Data Mining, Dedup & Audit Kualitas", "11")
# left chart
s.shapes.add_picture(os.path.join(CHARTS, "years.png"), Inches(0.5), y+Inches(0.1),
                     width=Inches(7.4))
# right stat cards
stats = [
    ("433", "Rekaman teridentifikasi", NAVY),
    ("144", "Duplikat dihapus (143 primer)", GREY),
    ("289", "Unik untuk screening", TEAL),
    ("0 / 0", "Missing DOI / Abstract", GREEN),
]
sx = Inches(8.2); sy = y+Inches(0.15); sw = Inches(4.55); shh = Inches(0.92)
for v, d, col in stats:
    card(s, sx, sy, sw, shh, fill=LGREY)
    rect(s, sx, sy, Inches(0.12), shh, fill=col)
    txt(s, sx+Inches(0.3), sy, Inches(1.7), shh, [(v, 26, col, True, False)],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, sx+Inches(2.0), sy, Inches(2.45), shh, [(d, 12.5, NAVY, False, False)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    sy += shh + Inches(0.12)
card(s, Inches(0.55), Inches(5.7), Inches(12.23), Inches(1.15), fill=NAVY)
txt(s, Inches(0.85), Inches(5.82), Inches(11.6), Inches(0.95),
    [[("Komposisi: ", 13, ACCENT, True, False),
      ("Article 169 · Journal Article 98 · Conference paper 87 · Case Report 2.   ", 12.5, WHITE, False, False),
      ("PICO-preview: 75% relevan → verdict PROCEED L3.", 12.5, ACCENT, True, False)]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s)

# =====================================================================
# SLIDE 12 — SCREENING & KAPPA
# =====================================================================
s = slide()
y = header(s, "Screening Dua-Reviewer (Modul 5)", "Reliabilitas (Cohen's Kappa) & PRISMA", "12")
s.shapes.add_picture(os.path.join(CHARTS, "kappa.png"), Inches(0.5), y+Inches(0.05),
                     width=Inches(7.3))
# PRISMA mini-flow on right
px = Inches(8.05); pw = Inches(4.7)
flow = [("289", "Disaring", TEAL), ("161", "Dieksklusi", RED), ("120", "Lolos ke full-text", GREEN)]
fy = y+Inches(0.1)
for i, (v, t, col) in enumerate(flow):
    card(s, px, fy, pw, Inches(0.92), fill=col)
    txt(s, px+Inches(0.25), fy, Inches(1.4), Inches(0.92), [(v, 26, WHITE, True, False)],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, px+Inches(1.7), fy, Inches(2.8), Inches(0.92), [(t, 14, WHITE, True, False)],
        anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow)-1:
        a = rect(s, px+pw/2-Inches(0.18), fy+Inches(0.92), Inches(0.36), Inches(0.2),
                 fill=GREY, shape=MSO_SHAPE.DOWN_ARROW)
    fy += Inches(1.12)
card(s, px, fy, pw, Inches(0.7), fill=LGREY, line=TEAL)
txt(s, px+Inches(0.2), fy, pw-Inches(0.4), Inches(0.7),
    [[("κ final = 0.71", 15, NAVY, True, False), ("  · Substantial Agreement", 12.5, GREY, False, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
# bottom note
txt(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.6),
    [[("6 iterasi kalibrasi (κ: 0.00 → 0.70) + 16 batch screening massal · 47 disagreement diselesaikan · ", 11.5, GREY, False, True),
      ("Eksklusi terbesar: I-NOMATCH 91 (bukan SSM/Mamba) & P-NOMATCH 54 (populasi).", 11.5, NAVY, True, False)]])
footer(s)

# =====================================================================
# SLIDE 13 — MODULE 6 (CURRENT) + QDRANT
# =====================================================================
s = slide()
y = header(s, "Posisi Saat Ini (Modul 6 · Langkah 1)", "Akuisisi Full-Text & Vektorisasi", "13")
chip(s, Inches(0.55), y, Inches(4.2), "STATUS: M6_STEP1 — WAITING_SYNC (QDRANT)", ACCENT, NAVY, 11)
stats = [
    ("120", "Paper masuk (include)", NAVY),
    ("118", "Full-text via HITL download", TEAL),
    ("2", "Open Access (Unpaywall/arXiv)", GREEN),
    ("0", "Inaccessible", GREY),
]
cw = Inches(2.95); gap = Inches(0.15); x0 = Inches(0.55); cy = y+Inches(0.65)
for i, (v, d, col) in enumerate(stats):
    x = x0 + i*(cw+gap)
    card(s, x, cy, cw, Inches(1.5), fill=LGREY)
    rect(s, x, cy, cw, Pt(5), fill=col)
    txt(s, x, cy+Inches(0.18), cw, Inches(0.7), [(v, 32, col, True, False)], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.15), cy+Inches(0.95), cw-Inches(0.3), Inches(0.5),
        [(d, 11.5, NAVY, False, False)], align=PP_ALIGN.CENTER, line_spacing=1.0)
# Qdrant box
card(s, Inches(0.55), Inches(4.5), Inches(6.0), Inches(2.35), fill=NAVY)
chip(s, Inches(0.8), Inches(4.7), Inches(3.0), "VECTOR DB · QDRANT", TEAL, WHITE, 11)
txt(s, Inches(0.8), Inches(5.25), Inches(5.5), Inches(1.5),
    [[("54", 40, ACCENT, True, False), ("  paper ter-vektorisasi", 16, WHITE, True, False)],
     [("2.918 chunks", 16, TEAL, True, False),
      (" tersimpan di koleksi ", 13, WHITE, False, False),
      ("scientific_articles", 13, ACCENT, False, True)],
     [("Siap untuk ekstraksi semantik (RAG) pada Modul 7.", 12.5, GREY, False, True)]],
    line_spacing=1.1, space_after=6)
# pipeline desc
card(s, Inches(6.75), Inches(4.5), Inches(6.03), Inches(2.35), fill=LGREY, line=TEAL)
txt(s, Inches(7.0), Inches(4.65), Inches(5.6), Inches(0.4),
    [("Alur Modul 6:", 14, NAVY, True, False)])
bullets(s, Inches(7.0), Inches(5.1), Inches(5.6), Inches(1.7), [
    "Resolusi OA otomatis (Unpaywall/arXiv) → unduh PDF.",
    "Human-in-the-loop untuk paper berbayar (118).",
    "Chunking + embedding → upsert ke Qdrant.",
    "Sinkronisasi sedang berjalan (langkah aktif).",
], size=12, gap=6, bullet_color=TEAL)
footer(s)

# =====================================================================
# SLIDE 14 — PROGRESS ROADMAP
# =====================================================================
s = slide()
y = header(s, "Status Kemajuan", "Peta Jalan Penelitian & Capaian", "14")
phases = [
    ("M1–M2", "Fondasi & PICO", "Topik, gap Tipe A, 4 RQ, FINER PASS", "done"),
    ("M3–M4", "Pencarian & Mining", "4 database, 433→289 rekaman unik", "done"),
    ("M5", "Screening", "κ=0.71, 120 paper lolos full-text", "done"),
    ("M6", "Akuisisi & Vektorisasi", "120 full-text · 54 di Qdrant (sinkronisasi)", "current"),
    ("M7", "Ekstraksi Data", "Ekstraksi metrik akurasi & efisiensi (RAG)", "next"),
    ("M8", "Sintesis & Bibliometrik", "Jawab RQ1–RQ4, analisis trade-off", "next"),
    ("M9", "Penyusunan Naskah", "Manuskrip tinjauan sistematis", "next"),
]
# progress bar
done = 5.5; total = 7
barx = Inches(0.55); bary = y; barw = Inches(12.23)
rect(s, barx, bary, barw, Inches(0.4), fill=LGREY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
prog = rect(s, barx, bary, Emu(int(barw* (done/total))), Inches(0.4), fill=TEAL,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, barx, bary, barw, Inches(0.4),
    [[("  ±78% pipeline SLR selesai  (M1–M5 tuntas, M6 berjalan)", 12, NAVY, True, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
ly = bary+Inches(0.7)
for tag, t, d, st in phases:
    rh = Inches(0.62)
    if st == "done":
        dot = GREEN; lab = "SELESAI"; labc = GREEN
    elif st == "current":
        dot = ACCENT; lab = "BERJALAN"; labc = ACCENT
    else:
        dot = GREY; lab = "RENCANA"; labc = GREY
    rect(s, barx+Inches(0.05), ly+Inches(0.1), Inches(0.4), Inches(0.4), fill=dot, shape=MSO_SHAPE.OVAL)
    txt(s, barx+Inches(0.65), ly, Inches(1.3), rh, [(tag, 14, NAVY, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, barx+Inches(2.0), ly, Inches(3.3), rh, [(t, 13.5, NAVY, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, barx+Inches(5.3), ly, Inches(5.6), rh, [(d, 12.5, GREY, False, False)], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, barx+Inches(11.0), ly+Inches(0.12), Inches(1.2), lab, labc, WHITE, 9.5)
    ly += rh + Inches(0.03)
footer(s)

# =====================================================================
# SLIDE 15 — CONTRIBUTION & NEXT STEPS
# =====================================================================
s = slide()
y = header(s, "Kontribusi & Rencana Lanjut", "Kontribusi yang Diharapkan & Langkah Berikutnya", "15")
card(s, Inches(0.55), y, Inches(6.0), Inches(4.5), fill=LGREY, line=TEAL)
chip(s, Inches(0.8), y+Inches(0.22), Inches(3.2), "KONTRIBUSI YANG DIHARAPKAN", TEAL, WHITE, 11)
bullets(s, Inches(0.85), y+Inches(0.85), Inches(5.5), Inches(3.5), [
    "Tinjauan sistematis pertama khusus arsitektur SSM/Mamba untuk neural decoding BCI.",
    "Sintesis kuantitatif trade-off akurasi ↔ latensi (<10 ms) lintas modalitas.",
    "Peta konfigurasi arsitektur (backbone murni vs hibrida POSSM) & keunggulannya.",
    "Panduan kelayakan model 'mangkus' untuk neuro-prostetik real-time di edge.",
    "Metodologi SLR berbantu agentic-AI yang reprodusibel & dapat diaudit.",
], size=13, gap=10, bullet_color=TEAL)
card(s, Inches(6.75), y, Inches(6.03), Inches(4.5), fill=NAVY)
chip(s, Inches(7.0), y+Inches(0.22), Inches(3.0), "LANGKAH BERIKUTNYA", ACCENT, NAVY, 11)
nexts = [
    ("M6 – selesaikan sync", "Finalisasi vektorisasi 120 paper ke Qdrant."),
    ("M7 – Ekstraksi", "Ekstraksi metrik akurasi & efisiensi via RAG."),
    ("M8 – Sintesis", "Jawab RQ1–RQ4 + analisis bibliometrik."),
    ("M9 – Naskah", "Susun manuskrip & sidang hasil."),
]
ny = y+Inches(0.9)
for t, d in nexts:
    rect(s, Inches(7.0), ny+Inches(0.05), Inches(0.12), Inches(0.75), fill=ACCENT)
    txt(s, Inches(7.25), ny, Inches(5.3), Inches(0.85),
        [[(t, 13.5, ACCENT, True, False)], [(d, 12, WHITE, False, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=2)
    ny += Inches(0.92)
footer(s)

# =====================================================================
# SLIDE 16 — CLOSING / THANK YOU
# =====================================================================
s = slide()
bg(s, NAVY)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(4.9), SW, Inches(2.6), fill=NAVY2)
rect(s, 0, Inches(4.9), SW, Pt(3), fill=TEAL)
random.seed(11)
for _ in range(22):
    cx = Inches(0.3 + random.random()*3.4)
    cy = Inches(0.3 + random.random()*2.4)
    d = Pt(4 + random.random()*7)
    rect(s, cx, cy, d, d, fill=(TEAL if random.random() > 0.5 else ACCENT), shape=MSO_SHAPE.OVAL)
txt(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.6),
    [[("Terima Kasih", 50, WHITE, True, False)],
     [("Diskusi & masukan sangat diharapkan untuk penyempurnaan proposal.", 17, GREY, False, True)]],
    space_after=14)
txt(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(1.8),
    [[("Rolly Maulana Awangga", 20, WHITE, True, False),
      ("   ·   NIM 303012510004", 14, GREY, False, False)],
     [("rollymaulanaa@student.telkomuniversity.ac.id  ·  Program Doktor — Universitas Telkom",
       14, TEAL, False, False)],
     [("Peningkatan Akurasi & Skalabilitas Neural Decoding via State-Space Models pada Ekosistem BCI",
       12.5, GREY, False, True)]], space_after=6)

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides._sldIdLst))
